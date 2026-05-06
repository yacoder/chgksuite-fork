import copy
import functools
import math
import os
import re
import urllib.parse

import toml

from chgksuite.common import (
    log_wrap,
    optimize_ooxml_images,
    replace_escaped,
    tryint,
)
from chgksuite.composer.composer_common import (
    BaseExporter,
    backtick_replace,
    parseimg,
    remove_accents_standalone,
)
from chgksuite.composer.docx import (
    _HYPERLINK_SAFE_CHARS,
    _docx_font_name,
    _docx_font_spec,
    _select_font_faces,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt


_PPTX_HYPERLINK_COLOR = (0x05, 0x63, 0xC1)
_EMU_PER_INCH = 914400
_PX_PER_INCH = 96
_PT_PER_INCH = 72
_PPTX_URL_BREAK_AFTER_CHARS = "/.-_\u2011?&=%#:"
_PPTX_URL_RE = re.compile(r"https?://|www\.")


@functools.lru_cache(maxsize=256)
def _pillow_font(font_path, pixel_size):
    from PIL import ImageFont

    return ImageFont.truetype(font_path, size=pixel_size)


def _optimize_size_enabled(args):
    return (getattr(args, "optimize_size", None) or "on") == "on"


def optimize_pptx_images(pptx_path, quality=80):
    return optimize_ooxml_images(
        pptx_path, media_prefix="ppt/media/", rels_prefix="ppt/", quality=quality
    )


def _get_pptx_handout_text_space_after(handout_cfg):
    return handout_cfg.get("text_space_after", handout_cfg.get("space_after", 18))


class PptxExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.config_path = os.path.abspath(self.args.pptx_config)
        with open(self.config_path, encoding="utf8") as f:
            self.c = toml.load(f)
        self.font_spec = _docx_font_spec(self.args)
        self.font_faces = None
        self.optimize_size = _optimize_size_enabled(self.args)
        self.font_name = _docx_font_name(self.font_spec)
        self._measurement_font_faces = None
        if self.font_name:
            font_cfg = self.c.setdefault("font", {})
            font_cfg["name"] = self.font_name
            if "heading_name" in font_cfg:
                font_cfg["heading_name"] = self.font_name
        self.qcount = 0
        hs = self.labels["question_labels"]["handout"]
        self.re_handout_1 = re.compile(
            "\\[" + hs + ".(?P<body>.+?)\\]", flags=re.DOTALL
        )
        self.re_handout_2 = re.compile("^" + hs + ".(?P<body>.+?)$")

    def _get_heading_font_name(self):
        try:
            font_cfg = self.c.get("font", {})
            return font_cfg.get("heading_name") or font_cfg.get("name")
        except Exception:
            return None

    def _get_font_size(self, key, fallback):
        font_cfg = self.c.get("font", {})
        if font_cfg.get(key) is not None:
            return font_cfg[key]
        if key == "question_size":
            if self.c.get("force_text_size_question") is not None:
                return self.c["force_text_size_question"]
            return self._get_font_size("default_size", fallback)
        if key == "answer_size":
            if self.c.get("force_text_size_answer") is not None:
                return self.c["force_text_size_answer"]
            return self._get_font_size("default_size", fallback)
        if key == "tour_size":
            return self._get_font_size("default_size", fallback)
        if key == "default_size":
            if self.c.get("force_text_size_question") is not None:
                return self.c["force_text_size_question"]
            text_size_grid = self.c.get("text_size_grid", {})
            if text_size_grid.get("default"):
                return text_size_grid["default"]
        return fallback

    def _get_grid_elements(self, role):
        text_size_grid = self.c.get("text_size_grid", {})
        return text_size_grid.get(f"{role}_elements") or text_size_grid.get(
            "elements", []
        )

    def _text_for_grid(self, text):
        if isinstance(text, list):
            return "\n".join(self._text_for_grid(element) for element in text)
        return str(text)

    def _get_grid_font_size(self, role, text, fallback):
        elements = self._get_grid_elements(role)
        if not elements:
            return fallback
        text_length = len(self._text_for_grid(text))
        for element in sorted(elements, key=lambda item: item["length"]):
            if text_length <= element["length"]:
                return element["size"]
        text_size_grid = self.c.get("text_size_grid", {})
        return text_size_grid.get(
            f"{role}_smallest", text_size_grid.get("smallest", fallback)
        )

    def _get_font_size_for_text(self, role, text, key, fallback):
        return self._get_grid_font_size(
            role, text, self._get_font_size(key, fallback)
        )

    def _apply_font_size_to_text_frame(self, text_frame, size, line_spacing_key=None):
        size = PptxPt(size)
        for p in text_frame.paragraphs:
            p.font.size = size
            self._set_line_spacing(p, size, line_spacing_key=line_spacing_key)
            for r in p.runs:
                r.font.size = size

    def _apply_font_to_text_frame(self, text_frame, font_name=None):
        if not font_name:
            return
        for p in text_frame.paragraphs:
            for r in p.runs:
                r.font.name = font_name

    def _alpha_marker(self, number, upper=False):
        result = ""
        while number:
            number, remainder = divmod(number - 1, 26)
            result = chr(ord("a") + remainder) + result
        if upper:
            return result.upper()
        return result

    def _roman_marker(self, number, upper=False):
        result = ""
        for value, numeral in (
            (1000, "m"),
            (900, "cm"),
            (500, "d"),
            (400, "cd"),
            (100, "c"),
            (90, "xc"),
            (50, "l"),
            (40, "xl"),
            (10, "x"),
            (9, "ix"),
            (5, "v"),
            (4, "iv"),
            (1, "i"),
        ):
            while number >= value:
                result += numeral
                number -= value
        if upper:
            return result.upper()
        return result

    def _format_list_marker(self, number):
        style = self.c.get("list", {}).get("numbering_style", "1.")
        if "{n}" in style:
            return style.format(n=number)
        if not style:
            style = "1."
        marker_type, suffix = style[0], style[1:]
        if marker_type == "1":
            marker = str(number)
        elif marker_type == "a":
            marker = self._alpha_marker(number)
        elif marker_type == "A":
            marker = self._alpha_marker(number, upper=True)
        elif marker_type == "i":
            marker = self._roman_marker(number)
        elif marker_type == "I":
            marker = self._roman_marker(number, upper=True)
        else:
            marker = str(number)
            suffix = style
        return f"{marker}{suffix}"

    def _include_handout_label(self):
        return bool(self.c.get("handout", {}).get("include_label", False))

    def _add_handout_on_separate_slide(self):
        add_handout_on_separate_slide = self.c.get("add_handout_on_separate_slide")
        return add_handout_on_separate_slide is None or add_handout_on_separate_slide

    def _disable_shrink_fit(self):
        return bool(self.c.get("disable_shrink_fit"))

    def _overlay_image_and_text(self):
        return bool(self.c.get("overlay_image_and_text"))

    def _service_slides_config(self):
        return self.c.get("service_slides", {})

    def _skip_generated_title_slide(self):
        return bool(self._service_slides_config().get("skip_generated_title"))

    def _slide_indices_from_config(self, key):
        value = self._service_slides_config().get(key)
        if value is None:
            return []
        if isinstance(value, list):
            return [int(index) for index in value]
        return [int(value)]

    def _configured_service_slide_indices(self):
        indices = []
        for key in (
            "intro",
            "between_tours",
            "final",
            "numbered_tours_stubs",
            "numbered_tour_stubs",
            "remove",
        ):
            indices.extend(self._slide_indices_from_config(key))
        return indices

    def _numbered_tour_stub_indices(self):
        return self._slide_indices_from_config(
            "numbered_tours_stubs"
        ) or self._slide_indices_from_config("numbered_tour_stubs")

    def _remap_relationship_ids(self, element, rel_id_map):
        if not rel_id_map:
            return
        for child in element.iter():
            for attr_name, attr_value in child.attrib.items():
                if attr_value in rel_id_map:
                    child.set(attr_name, rel_id_map[attr_value])

    def _copy_slide_background(self, source_slide, slide):
        source_bg = source_slide.element.cSld.bg
        if source_bg is None:
            return
        destination_bg = slide.element.cSld.bg
        if destination_bg is not None:
            slide.element.cSld.remove(destination_bg)
        slide.element.cSld.insert(0, copy.deepcopy(source_bg))

    def _copy_slide_relationships(self, source_slide, slide):
        rel_id_map = {}
        for rel in source_slide.part.rels.values():
            if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith(
                "/notesSlide"
            ):
                continue
            target = rel._target if rel.is_external else rel.target_part
            rel_id_map[rel.rId] = slide.part.rels._add_relationship(
                rel.reltype, target, rel.is_external
            )
        return rel_id_map

    def _clone_slide(self, source_slide):
        slide = self.prs.slides.add_slide(source_slide.slide_layout)
        for shape in list(slide.shapes):
            self._remove_shape(shape)

        self._copy_slide_background(source_slide, slide)
        rel_id_map = self._copy_slide_relationships(source_slide, slide)
        for shape in source_slide.shapes:
            element = copy.deepcopy(shape.element)
            self._remap_relationship_ids(element, rel_id_map)
            slide.shapes._spTree.insert_element_before(element, "p:extLst")
        return slide

    def _remove_slide_at(self, index):
        slide_id = self.prs.slides._sldIdLst[index]
        self.prs.slides._sldIdLst.remove(slide_id)
        self.prs.part.drop_rel(slide_id.rId)

    def _prepare_service_slide_templates(self):
        self._service_slide_templates = {}
        self._service_slide_indices_to_remove = []
        configured_indices = self._configured_service_slide_indices()
        if not configured_indices:
            return

        slide_count = len(self.prs.slides)
        for index in configured_indices:
            if index < 0 or index >= slide_count:
                raise ValueError(
                    f"service slide index {index} is out of range for "
                    f"{self.c['template_path']}"
                )

        for key in ("intro", "between_tours", "final"):
            self._service_slide_templates[key] = [
                self.prs.slides[index]
                for index in self._slide_indices_from_config(key)
            ]
        self._service_slide_templates["numbered_tours_stubs"] = [
            self.prs.slides[index] for index in self._numbered_tour_stub_indices()
        ]
        self._service_slide_indices_to_remove = sorted(
            set(configured_indices), reverse=True
        )

    def _remove_service_slide_templates(self):
        for index in self._service_slide_indices_to_remove:
            self._remove_slide_at(index)

    def _add_service_slides(self, key):
        for slide in getattr(self, "_service_slide_templates", {}).get(key, []):
            self._clone_slide(slide)

    def _add_numbered_tour_stub(self):
        slides = getattr(self, "_service_slide_templates", {}).get(
            "numbered_tours_stubs", []
        )
        tour_index = getattr(self, "_processed_tour_count", 0)
        if tour_index < len(slides):
            self._clone_slide(slides[tour_index])

    def _should_add_between_tours_slide(self, buffer):
        if not getattr(self, "_processed_question_count", 0):
            return False
        return any(element[0] == "section" for element in buffer)

    def _line_spacing_configured(self):
        font_cfg = self.c.get("font", {})
        if font_cfg.get("fixed_line_spacing"):
            return True
        if font_cfg.get("line_spacing_multiplier") is not None:
            return True
        return any(
            key.startswith("fixed_line_spacing_") and value is not None
            for key, value in font_cfg.items()
        )

    def _get_fixed_line_spacing(self, line_spacing_key):
        if not line_spacing_key:
            return None
        value = self.c.get("font", {}).get(f"fixed_line_spacing_{line_spacing_key}")
        if value is None:
            return None
        return PptxPt(value)

    def _set_line_spacing(self, paragraph, font_size, line_spacing_key=None):
        fixed_line_spacing = self._get_fixed_line_spacing(line_spacing_key)
        if fixed_line_spacing is not None:
            paragraph.line_spacing = fixed_line_spacing
            return
        font_cfg = self.c.get("font", {})
        multiplier = font_cfg.get("line_spacing_multiplier")
        if multiplier is not None:
            paragraph.line_spacing = float(multiplier)
            return
        if font_cfg.get("fixed_line_spacing"):
            paragraph.line_spacing = font_size

    def _set_paragraph_alignment(self, paragraph, align):
        if not align:
            return
        paragraph.alignment = getattr(PP_ALIGN, align.upper())

    def _configure_paragraph(
        self, paragraph, size=None, align=None, line_spacing_key=None
    ):
        paragraph.font.name = self.c["font"]["name"]
        font_size = PptxPt(size or self._get_font_size("default_size", 32))
        paragraph.font.size = font_size
        self._set_line_spacing(
            paragraph, font_size, line_spacing_key=line_spacing_key
        )
        self._set_paragraph_alignment(paragraph, align)
        return paragraph

    def _get_handout_space_after(self):
        return self.c.get("handout", {}).get("space_after", 18)

    def _get_handout_text_space_after(self):
        return _get_pptx_handout_text_space_after(self.c.get("handout", {}))

    def _get_handout_font_size(self):
        handout_cfg = self.c.get("handout", {})
        if handout_cfg.get("font_size") is not None:
            return handout_cfg["font_size"]
        return self._get_font_size("tour_size", 42)

    def _get_measurement_font_path(self):
        faces = self._get_measurement_font_faces()
        face = faces.get("regular")
        return getattr(face, "path", None) if face else None

    def _get_measurement_font_faces(self):
        if self.font_faces:
            return self.font_faces

        font_spec = self.font_spec or self.c.get("font", {}).get("name")
        if not font_spec:
            return {}
        if self._measurement_font_faces is not None:
            return self._measurement_font_faces
        try:
            self._measurement_font_faces = _select_font_faces(font_spec)
        except (OSError, PermissionError, ValueError):
            self._measurement_font_faces = {}
        return self._measurement_font_faces

    def _get_measurement_font_path_for_style(self, bold=False, italic=False):
        faces = self._get_measurement_font_faces()
        if not faces:
            return None
        if bold and italic:
            face = faces.get("bold_italic") or faces.get("bold") or faces.get("italic")
        elif bold:
            face = faces.get("bold")
        elif italic:
            face = faces.get("italic")
        else:
            face = faces.get("regular")
        face = face or faces.get("regular") or next(iter(faces.values()))
        return getattr(face, "path", None) if face else None

    def _font_pixel_size(self, font_size):
        return max(1, round(float(font_size) * _PX_PER_INCH / _PT_PER_INCH))

    def _measure_text_width_px(self, font_path, text, font_size):
        font = _pillow_font(font_path, self._font_pixel_size(font_size))
        if hasattr(font, "getlength"):
            return font.getlength(text)
        bbox = font.getbbox(text)
        return bbox[2] - bbox[0]

    def _measure_line_height_px(self, font_path, font_size):
        font = _pillow_font(font_path, self._font_pixel_size(font_size))
        height = getattr(getattr(font, "font", None), "height", None)
        if height:
            return height
        ascent, descent = font.getmetrics()
        return ascent + descent

    def _textbox_inner_width_px(self, textbox):
        text_frame = textbox.text_frame
        margin_left = text_frame.margin_left or 0
        margin_right = text_frame.margin_right or 0
        inner_width = max(textbox.width - margin_left - margin_right, 1)
        return inner_width / _EMU_PER_INCH * _PX_PER_INCH

    def _textbox_inner_height_px(self, textbox):
        text_frame = textbox.text_frame
        margin_top = text_frame.margin_top or 0
        margin_bottom = text_frame.margin_bottom or 0
        inner_height = max(textbox.height - margin_top - margin_bottom, 1)
        return inner_height / _EMU_PER_INCH * _PX_PER_INCH

    def _pt_to_px(self, value):
        return float(value) * _PX_PER_INCH / _PT_PER_INCH

    def _length_to_px(self, value):
        if value is None:
            return 0
        return value / _EMU_PER_INCH * _PX_PER_INCH

    def _effective_paragraph_size(self, paragraph):
        if paragraph.font.size:
            return paragraph.font.size.pt
        return self._get_font_size("default_size", 32)

    def _effective_run_size(self, run, paragraph):
        if run.font.size:
            return run.font.size.pt
        return self._effective_paragraph_size(paragraph)

    def _run_text_width_px(self, run, paragraph, text=None):
        text = run.text if text is None else text
        if not text:
            return 0
        font_path = self._get_measurement_font_path_for_style(
            bold=bool(run.font.bold), italic=bool(run.font.italic)
        )
        if not font_path:
            return None
        return self._measure_text_width_px(
            font_path, text, self._effective_run_size(run, paragraph)
        )

    def _run_line_height_px(self, run, paragraph):
        font_path = self._get_measurement_font_path_for_style(
            bold=bool(run.font.bold), italic=bool(run.font.italic)
        )
        if not font_path:
            return self._pt_to_px(self._effective_run_size(run, paragraph))
        return self._measure_line_height_px(
            font_path, self._effective_run_size(run, paragraph)
        )

    def _text_width_px(self, text, size, bold=False, italic=False):
        if not text:
            return 0
        font_path = self._get_measurement_font_path_for_style(
            bold=bold, italic=italic
        )
        if not font_path:
            return None
        return self._measure_text_width_px(font_path, text, size)

    def _token_width_px(self, paragraph, runs):
        total = 0
        for run, text in runs:
            width = self._run_text_width_px(run, paragraph, text=text)
            if width is None:
                return None
            total += width
        return total

    def _run_can_break_like_url(self, run, text):
        hyperlink = getattr(run, "hyperlink", None)
        if hyperlink is not None and getattr(hyperlink, "address", None):
            return True
        return bool(_PPTX_URL_RE.search(text))

    def _append_url_wrapping_tokens(self, line, run, text):
        token = ""
        for char in text:
            token += char
            if char in _PPTX_URL_BREAK_AFTER_CHARS:
                line.append((run, token, False))
                line.append((run, "", True))
                token = ""
        if token:
            line.append((run, token, False))

    def _split_runs_for_wrapping(self, paragraph):
        lines = [[]]
        runs_by_element = {id(run._r): run for run in paragraph.runs}

        def append_text(run, text):
            for line_index, line in enumerate(re.split(r"[\n\r\v]", text)):
                if line_index:
                    lines.append([])
                parts = re.split(r"( +)", line)
                for part in parts:
                    if not part:
                        continue
                    if part.isspace() and "\u00a0" not in part and "\u2011" not in part:
                        lines[-1].append((run, part, True))
                    elif self._run_can_break_like_url(run, part):
                        self._append_url_wrapping_tokens(lines[-1], run, part)
                    else:
                        lines[-1].append((run, part, False))

        for child in paragraph._p:
            if child.tag.endswith("}br"):
                lines.append([])
                continue
            if not child.tag.endswith("}r"):
                continue
            run = runs_by_element.get(id(child))
            if run is None:
                continue
            append_text(run, run.text)
        return lines

    def _paragraph_layout_lines(self, paragraph, max_width):
        lines = []
        max_unbreakable_width = 0
        for explicit_line in self._split_runs_for_wrapping(paragraph):
            if not explicit_line:
                lines.append([])
                continue

            current_line = []
            current_width = 0
            pending_spaces = []
            has_content = False

            for run, text, breakable_space in explicit_line:
                token_width = self._token_width_px(paragraph, [(run, text)])
                if token_width is None:
                    return None
                if breakable_space:
                    pending_spaces.append((run, text))
                    continue

                space_width = self._token_width_px(paragraph, pending_spaces)
                if space_width is None:
                    return None
                candidate_width = current_width + space_width + token_width
                max_unbreakable_width = max(max_unbreakable_width, token_width)
                if has_content and candidate_width > max_width:
                    lines.append(current_line)
                    current_line = [(run, text)]
                    current_width = token_width
                else:
                    current_line.extend(pending_spaces)
                    current_line.append((run, text))
                    current_width = candidate_width
                pending_spaces = []
                has_content = True

            lines.append(current_line)

        return lines, max_unbreakable_width

    def _paragraph_line_count(self, paragraph, max_width):
        layout = self._paragraph_layout_lines(paragraph, max_width)
        if layout is None:
            return None
        lines, max_unbreakable_width = layout
        return len(lines), max_unbreakable_width

    def _paragraph_line_height_px(self, paragraph, line_runs=None):
        text_runs = [(run, text) for run, text in line_runs or [] if text]
        if text_runs:
            line_height = max(
                self._run_line_height_px(run, paragraph) for run, _text in text_runs
            )
            font_size = max(
                self._effective_run_size(run, paragraph) for run, _text in text_runs
            )
        else:
            font_size = self._effective_paragraph_size(paragraph)
            font_path = self._get_measurement_font_path_for_style()
            if font_path:
                line_height = self._measure_line_height_px(font_path, font_size)
            else:
                line_height = self._pt_to_px(font_size)
        line_spacing = paragraph.line_spacing
        if line_spacing is None:
            return line_height
        if isinstance(line_spacing, float):
            return line_height * line_spacing
        return self._length_to_px(line_spacing)

    def _text_frame_fits(self, textbox):
        max_width = self._textbox_inner_width_px(textbox) * 0.99
        max_height = self._textbox_inner_height_px(textbox) * 0.99
        total_height = 0
        for paragraph in textbox.text_frame.paragraphs:
            line_info = self._paragraph_layout_lines(paragraph, max_width)
            if line_info is None:
                return True
            lines, max_unbreakable_width = line_info
            if max_unbreakable_width > max_width:
                return False
            for line_runs in lines:
                total_height += self._paragraph_line_height_px(paragraph, line_runs)
            total_height += self._length_to_px(paragraph.space_after)
            total_height += self._length_to_px(paragraph.space_before)
        return total_height <= max_height

    def _collect_textbox_font_sizes(self, textbox):
        items = []
        fallback = self._get_font_size("default_size", 32)
        for paragraph in textbox.text_frame.paragraphs:
            paragraph_size = paragraph.font.size.pt if paragraph.font.size else fallback
            items.append((paragraph.font, paragraph_size))
            for run in paragraph.runs:
                run_size = run.font.size.pt if run.font.size else paragraph_size
                items.append((run.font, run_size))
        return items

    def _set_textbox_font_delta(self, items, delta, min_size):
        for font, original_size in items:
            font.size = PptxPt(max(float(min_size), float(original_size) - delta))

    def _custom_shrink_textbox(self, textbox, min_size=None):
        if self._disable_shrink_fit():
            return
        if not self._get_measurement_font_faces():
            return
        items = self._collect_textbox_font_sizes(textbox)
        if not items:
            return
        min_size = min_size or self.c.get("text_size_grid", {}).get("smallest", 14)
        max_delta = max(
            0, math.ceil(max(original_size for _font, original_size in items) - min_size)
        )
        for delta in range(max_delta + 1):
            self._set_textbox_font_delta(items, delta, min_size)
            if self._text_frame_fits(textbox):
                return

    def _plain_text_lines_for_measurement(self, text):
        lines = [""]

        def append_text(value):
            parts = str(value).split("\n")
            lines[-1] += parts[0]
            for part in parts[1:]:
                lines.append(part)

        for run in self.parse_4s_elem(backtick_replace(text)):
            if run[0] == "screen":
                append_text(run[1]["for_screen"])
            elif run[0] == "linebreak":
                lines.append("")
            elif run[0] in ("img", "pagebreak"):
                continue
            else:
                append_text(run[1])
        return [line for line in lines if line.strip()]

    def _get_handout_font_size_for_text(self, text, textbox, min_size=None):
        max_size = self._get_handout_font_size()
        min_size = min_size or self._get_font_size("question_size", 32)
        if max_size <= min_size:
            return max_size

        font_path = self._get_measurement_font_path()
        if not font_path:
            return max_size

        lines = self._plain_text_lines_for_measurement(text)
        if not lines:
            return max_size

        max_width = self._textbox_inner_width_px(textbox) * 0.99
        current_size = float(max_size)
        min_size = float(min_size)
        while current_size > min_size:
            if all(
                self._measure_text_width_px(font_path, line, current_size) <= max_width
                for line in lines
            ):
                return current_size
            current_size -= 1
        return min_size

    def _get_handout_image_scale(self):
        return self.c.get("handout", {}).get("image_scale", 1)

    def _format_links_enabled(self):
        return self.c.get("format_links", True)

    def _get_image_space_after(self, image):
        if image.get("handout"):
            return PptxPt(self._get_handout_space_after())
        return 0

    def _scale_image_for_pptx(self, image):
        if not image.get("handout"):
            return image
        scale = self._get_handout_image_scale()
        if not scale or scale == 1:
            return image
        image["width"] *= scale
        image["height"] *= scale
        return image

    def _prepare_text_frame(self, text_frame):
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

    def _remove_shape(self, shape):
        element = shape._element
        element.getparent().remove(element)

    def get_textbox_qnumber(self, slide):
        kwargs = {}
        for param in ("left", "top", "width", "height"):
            try:
                kwargs[param] = PptxInches(self.c["number_textbox"][param])
            except KeyError:
                pass

        return self.get_textbox(slide, **kwargs)

    def get_textbox(self, slide, left=None, top=None, width=None, height=None):
        if left is None:
            left = PptxInches(self.c["textbox"]["left"])
        if top is None:
            top = PptxInches(self.c["textbox"]["top"])
        if width is None:
            width = PptxInches(self.c["textbox"]["width"])
        if height is None:
            height = PptxInches(self.c["textbox"]["height"])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._prepare_text_frame(textbox.text_frame)
        return textbox

    def _apply_run_defaults(self, run, para, color=None):
        if para.font.name:
            run.font.name = para.font.name
        if para.font.size:
            run.font.size = para.font.size
        if color is None:
            color = self.c["textbox"].get("color")
        if color:
            run.font.color.rgb = RGBColor(*color)
        if self.args.language == "ru":
            run.font.language_id = MSO_LANGUAGE_ID.RUSSIAN

    def add_runs(self, para, text, color=None):
        runs = []
        for index, part in enumerate(str(text).split("\n")):
            if index:
                para._p.append(OxmlElement("a:br"))
            if not part:
                continue
            run = para.add_run()
            run.text = part
            self._apply_run_defaults(run, para, color=color)
            runs.append(run)
        if not runs:
            run = para.add_run()
            run.text = ""
            self._apply_run_defaults(run, para, color=color)
            runs.append(run)
        return runs

    def add_run(self, para, text, color=None):
        return self.add_runs(para, text, color=color)[-1]

    def _apply_text_style(self, runs, style):
        for run in runs:
            if "italic" in style:
                run.font.italic = True
            if "bold" in style:
                run.font.bold = True
            if "underline" in style:
                run.font.underline = True

    def _add_styled_runs(self, para, text, style="", color=None):
        runs = self.add_runs(para, text, color=color)
        self._apply_text_style(runs, style)
        return runs

    def add_hyperlink_runs(self, para, text, url):
        if not self._format_links_enabled():
            return self.add_runs(para, text)
        hyperlink_url = urllib.parse.quote(url, safe=_HYPERLINK_SAFE_CHARS)
        runs = self.add_runs(para, text)
        for run in runs:
            run.hyperlink.address = hyperlink_url
            run.font.underline = True
            run.font.color.rgb = RGBColor(*_PPTX_HYPERLINK_COLOR)
        return runs

    def pptx_format(
        self, el, para, tf, slide, replace_spaces=True, blank_lines_between_items=False
    ):
        def r_sp(text):
            if replace_spaces:
                return self._replace_no_break(text)
            return text

        if isinstance(el, list):
            if len(el) > 1 and isinstance(el[1], list):
                self.pptx_format(
                    el[0],
                    para,
                    tf,
                    slide,
                    blank_lines_between_items=blank_lines_between_items,
                )
                blank_line = self.c.get("list", {}).get(
                    "blank_line_before_items", True
                )
                for licount, li in enumerate(el[1], start=1):
                    if blank_line and (licount == 1 or blank_lines_between_items):
                        prefix = "\n\n"
                    else:
                        prefix = "\n"
                    marker = self._format_list_marker(licount)
                    self.add_run(para, f"{prefix}{marker} ")
                    self.pptx_format(
                        li,
                        para,
                        tf,
                        slide,
                        blank_lines_between_items=blank_lines_between_items,
                    )
            else:
                blank_line = self.c.get("list", {}).get(
                    "blank_line_before_items", True
                )
                for licount, li in enumerate(el, start=1):
                    if blank_line and blank_lines_between_items and licount > 1:
                        prefix = "\n\n"
                    else:
                        prefix = "\n"
                    marker = self._format_list_marker(licount)
                    self.add_run(para, f"{prefix}{marker} ")
                    self.pptx_format(
                        li,
                        para,
                        tf,
                        slide,
                        blank_lines_between_items=blank_lines_between_items,
                    )

        if isinstance(el, str):
            self.logger.debug("parsing element {}:".format(log_wrap(el)))
            el = backtick_replace(el)

            for run in self.parse_4s_elem(el):
                if run[0] == "screen":
                    self.add_runs(para, r_sp(run[1]["for_screen"]))

                elif run[0] == "linebreak":
                    self.add_run(para, "\n")

                elif run[0] == "strike":
                    runs = self.add_runs(para, r_sp(run[1]))
                    for r in runs:
                        r.font.strike = True  # TODO: doesn't work as of 2023-12-24, cf. https://github.com/scanny/python-pptx/issues/339

                elif run[0] == "hyperlink":
                    self.add_hyperlink_runs(para, run[1], run[1])

                elif run[0] == "img":
                    pass  # image processing is moved to other places

                else:
                    self._add_styled_runs(para, r_sp(run[1]), run[0])

    def pptx_process_text(
        self,
        s,
        image=None,
        strip_brackets=True,
        replace_spaces=True,
        do_not_remove_accents=False,
    ):
        hs = self.regexes["handout_short"]
        if isinstance(s, list):
            for i in range(len(s)):
                s[i] = self.pptx_process_text(s[i], image=image)
            return s
        if not (self.args.do_not_remove_accents or do_not_remove_accents):
            s = remove_accents_standalone(s, self.regexes)
        if strip_brackets:
            s = self.remove_square_brackets(s)
            s = s.replace("]\n", "]\n\n")
        else:
            s = replace_escaped(s)
        if image:
            s = re.sub("\\[" + hs + "(.+?)\\]", "", s, flags=re.DOTALL)
            s = s.strip()
        elif re.search(hs, s) and not self._include_handout_label():
            re_hs = re.search("\\[" + hs + ".+?: ?(.+)\\]", s, flags=re.DOTALL)
            if re_hs:
                s = s.replace(re_hs.group(0), re_hs.group(1))
        s = re.sub(" +", " ", s)
        for punct in (".", ",", "!", "?", ":"):
            s = s.replace(" " + punct, punct)
        if replace_spaces:
            s = self._replace_no_break(s)
        s = s.strip()
        return s

    def apply_vertical_alignment_if_needed(self, text_frame):
        self._prepare_text_frame(text_frame)
        align = self.c["textbox"].get("vertical_align")
        if align:
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.vertical_anchor = getattr(MSO_VERTICAL_ANCHOR, align.upper())

    def _get_title_textbox_dimension(self, key, fallback):
        title_cfg = self.c.get("title_textbox", {})
        if key in title_cfg:
            return PptxInches(title_cfg[key])
        textbox_cfg = self.c.get("textbox", {})
        if key in textbox_cfg:
            return PptxInches(textbox_cfg[key])
        return fallback

    def format_title_slide(self, title, subtitle=None):
        if title is None or not hasattr(title, "text_frame"):
            return
        tf = title.text_frame
        self._prepare_text_frame(tf)
        self._apply_font_size_to_text_frame(
            tf, self._get_font_size("title_size", 60), line_spacing_key="title"
        )
        if subtitle is None:
            layout_title = None
            try:
                title_idx = title.placeholder_format.idx
                for layout_shape in self.TITLE_SLIDE.shapes:
                    if layout_shape.placeholder_format.idx == title_idx:
                        layout_title = layout_shape
                        break
            except (AttributeError, ValueError):
                pass
            default_left = layout_title.left if layout_title else PptxInches(1.67)
            default_width = (
                layout_title.width
                if layout_title
                else self.prs.slide_width - 2 * default_left
            )
            title.left = self._get_title_textbox_dimension("left", default_left)
            title.width = self._get_title_textbox_dimension("width", default_width)
            title.top = self._get_title_textbox_dimension(
                "top", PptxInches(0.8)
            )
            title.height = self._get_title_textbox_dimension(
                "height", PptxInches(6.1)
            )
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            self._custom_shrink_textbox(title)
        elif hasattr(subtitle, "text_frame"):
            self._prepare_text_frame(subtitle.text_frame)
            self._apply_font_size_to_text_frame(
                subtitle.text_frame,
                self._get_font_size("default_size", 32),
                line_spacing_key="default",
            )
            self._custom_shrink_textbox(subtitle)

    def _process_block(self, block):
        section = [x for x in block if x[0] == "section"]
        editor = [x for x in block if x[0] == "editor"]
        meta = [x for x in block if x[0] == "meta"]
        if not section and not editor and not meta:
            return
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True
        p = self.init_paragraph(tf)
        add_line_break = False
        if section:
            if self.c.get("tour_as_question_number"):
                txt = self.pptx_process_text(section[0][1])
                if self.c.get("tour_as_question_number") == "caps":
                    txt = txt.upper()
                self.set_question_number(slide, number=txt)
            else:
                r = self.add_run(
                    p, self._replace_no_break(self.pptx_process_text(section[0][1]))
                )
                heading_font = self._get_heading_font_name()
                if heading_font:
                    r.font.name = heading_font
                r.font.size = PptxPt(self._get_font_size("tour_size", 32))
                add_line_break = True
        if editor:
            if add_line_break:
                self.add_run(p, "\n\n")
            self.pptx_format(
                self.pptx_process_text(editor[0][1]),
                p,
                tf,
                slide,
            )
            add_line_break = True
        if meta:
            for element in meta:
                if add_line_break:
                    self.add_run(p, "\n\n")
                self.pptx_format(
                    self.pptx_process_text(element[1]),
                    p,
                    tf,
                    slide,
                )
                add_line_break = True
        self._custom_shrink_textbox(textbox)

    def process_buffer(self, buffer):
        heading_block = []
        editor_block = []
        section_block = []
        block = heading_block
        for element in buffer:
            if element[0] == "section":
                block = section_block
            if element[0] == "editor" and not section_block:
                block = editor_block
            block.append(element)
        heading = [x for x in heading_block if x[0] == "heading"]
        ljheading = [x for x in heading_block if x[0] == "ljheading"]
        title_text = ljheading or heading
        date_text = [x for x in heading_block if x[0] == "date"]
        if title_text and not self._skip_generated_title_slide():
            if len(self.prs.slides) == 1 and not getattr(
                self, "_service_slide_templates", {}
            ):
                slide = self.prs.slides[0]
            else:
                slide = self.prs.slides.add_slide(self.TITLE_SLIDE)
            title = slide.shapes.title
            title.text = title_text[0][1]
            subtitle = None
            if date_text:
                try:
                    subtitle = slide.placeholders[1]
                    subtitle.text = date_text[0][1]
                except KeyError:
                    pass
            else:
                try:
                    self._remove_shape(slide.placeholders[1])
                except KeyError:
                    pass
            self.format_title_slide(title, subtitle=subtitle)
            heading_font = self._get_heading_font_name()
            if heading_font:
                if title is not None and hasattr(title, "text_frame"):
                    self._apply_font_to_text_frame(title.text_frame, heading_font)
                if subtitle is not None and hasattr(subtitle, "text_frame"):
                    self._apply_font_to_text_frame(subtitle.text_frame, heading_font)
        self._process_block(editor_block)
        if section_block:
            self._add_numbered_tour_stub()
            self._process_block(section_block)
            self._processed_tour_count += 1

    def set_question_number(self, slide, number):
        if self.args.disable_numbers:
            return
        qntextbox = self.get_textbox_qnumber(slide)
        qtf = qntextbox.text_frame
        qtf_p = self.init_paragraph(qtf)
        if self.c["number_textbox"].get("align"):
            qtf_p.alignment = getattr(
                PP_ALIGN, self.c["number_textbox"]["align"].upper()
            )
        if (
            self.c.get("question_number_format") == "caps"
            and tryint(number) is not None
        ):
            number = f"ВОПРОС {number}"
        qtf_r = self.add_run(qtf_p, number)
        if self.c["number_textbox"].get("bold"):
            qtf_r.font.bold = True
        if self.c["number_textbox"].get("color"):
            qtf_r.font.color.rgb = RGBColor(*self.c["number_textbox"]["color"])
        number_font_size = self.c["number_textbox"].get("font_size")
        if number_font_size is None:
            number_font_size = self._get_font_size("number_size", None)
        if number_font_size is not None:
            number_font_size = PptxPt(number_font_size)
            qtf_r.font.size = number_font_size
            self._set_line_spacing(
                qtf_p, number_font_size, line_spacing_key="number"
            )

    def _get_handout_from_4s(self, text):
        if isinstance(text, list):
            for el in text:
                handout = self._get_handout_from_4s(el)
                if handout:
                    return handout
        elif isinstance(text, str):
            match_ = self.re_handout_1.search(text)
            if match_:
                if self._include_handout_label():
                    return match_.group(0)
                return match_.group("body")
            else:
                lines = text.split("\n")
                for line in lines:
                    match_ = self.re_handout_2.search(line)
                    if match_:
                        if self._include_handout_label():
                            return match_.group(0)
                        return match_.group("body")

    def _split_handout_from_text(self, text):
        if not isinstance(text, str):
            return None, text
        match_ = self.re_handout_1.search(text)
        if match_:
            if self._include_handout_label():
                handout = match_.group(0)
            else:
                handout = match_.group("body")
            question = f"{text[: match_.start()]}{text[match_.end() :]}"
            return handout.strip(), question.strip()

        lines = text.split("\n")
        for index, line in enumerate(lines):
            match_ = self.re_handout_2.search(line.strip())
            if not match_:
                continue
            if self._include_handout_label():
                handout = line
            else:
                handout = match_.group("body")
            question = "\n".join(lines[:index] + lines[index + 1 :])
            return handout.strip(), question.strip()

        return None, text

    def _get_image_from_4s(self, text):
        if isinstance(text, list):
            for el in text:
                image = self._get_image_from_4s(el)
                if image:
                    return image
        elif isinstance(text, str):
            handout_match = self.re_handout_1.search(text)
            for run in self.parse_4s_elem(text):
                if run[0] == "img":
                    parsed_image = parseimg(
                        run[1],
                        dimensions="inches",
                        tmp_dir=self.dir_kwargs.get("tmp_dir"),
                        targetdir=self.dir_kwargs.get("targetdir"),
                    )
                    parsed_image["handout"] = bool(
                        handout_match and run[1] in handout_match.group(0)
                    )
                    return self._scale_image_for_pptx(parsed_image)

    def make_slide_layout(self, image, slide, allowbigimage=True):
        if image:
            ratio = image["width"] / image["height"]
            img_base_width = PptxInches(image["width"])
            img_base_height = PptxInches(image["height"])
            base_left = PptxInches(self.c["textbox"]["left"])
            base_top = PptxInches(self.c["textbox"]["top"])
            base_width = PptxInches(self.c["textbox"]["width"])
            base_height = PptxInches(self.c["textbox"]["height"])
            image_space_after = self._get_image_space_after(image)
            if self._overlay_image_and_text():
                slide.shapes.add_picture(
                    image["imgfile"],
                    left=base_left,
                    top=base_top,
                    width=img_base_width,
                    height=img_base_height,
                )
                return self.get_textbox(slide)
            if self.c.get("disable_autolayout"):
                slide.shapes.add_picture(
                    image["imgfile"],
                    left=base_left,
                    top=base_top,
                    width=img_base_width,
                    height=img_base_height,
                )
                if ratio < 1:  # vertical image
                    left = base_left + img_base_width + image_space_after
                    top = base_top
                    width = max(base_width - img_base_width - image_space_after, 0)
                    height = base_height
                else:  # horizontal/square image
                    left = base_left
                    top = base_top + img_base_height + image_space_after
                    width = base_width
                    height = max(base_height - img_base_height - image_space_after, 0)
                return self.get_textbox(
                    slide, left=left, top=top, width=width, height=height
                )
            big_mode = (
                image["big"] and not self.c.get("text_is_duplicated") and allowbigimage
            )
            if ratio < 1:  # vertical image
                max_width = base_width // 3
                if big_mode:
                    max_width *= 2
                if image.get("handout"):
                    max_width = int(max_width * self._get_handout_image_scale())
                max_width = min(max_width, base_width - image_space_after)
                if img_base_width > max_width or big_mode:
                    img_width = max_width
                    img_height = int(img_base_height * (max_width / img_base_width))
                else:
                    img_width = img_base_width
                    img_height = img_base_height
                left = base_left + img_width + image_space_after
                top = base_top
                width = max(base_width - img_width - image_space_after, 0)
                height = base_height
                img_left = base_left
                img_top = int(base_top + 0.5 * (base_height - img_height))
            else:  # horizontal/square image
                max_height = base_height // 3
                if big_mode:
                    max_height *= 2
                if image.get("handout"):
                    max_height = int(max_height * self._get_handout_image_scale())
                max_height = min(max_height, base_height - image_space_after)
                if img_base_height > max_height or big_mode:
                    img_height = max_height
                    img_width = int(img_base_width * (max_height / img_base_height))
                else:
                    img_width = img_base_width
                    img_height = img_base_height
                left = base_left
                top = base_top + img_height + image_space_after
                width = base_width
                height = max(base_height - img_height - image_space_after, 0)
                img_top = base_top
                img_left = int(base_left + 0.5 * (base_width - img_width))
            slide.shapes.add_picture(
                image["imgfile"],
                left=img_left,
                top=img_top,
                width=img_width,
                height=img_height,
            )
            textbox = slide.shapes.add_textbox(left, top, width, height)
            self._prepare_text_frame(textbox.text_frame)
            return textbox
        else:
            return self.get_textbox(slide)

    def add_slide_with_image(self, image, number=None):
        slide = self.prs.slides.add_slide(self.QUESTION_SLIDE)
        if number:
            self.set_question_number(slide, number)
        img_width = PptxInches(image["width"])
        img_height = PptxInches(image["height"])
        base_left = PptxInches(self.c["textbox"]["left"])
        base_top = PptxInches(self.c["textbox"]["top"])
        base_width = PptxInches(self.c["textbox"]["width"])
        base_height = PptxInches(self.c["textbox"]["height"])
        if image["big"] or img_width > base_width:
            img_width, img_height = (
                base_width,
                int(img_height * (base_width / img_width)),
            )
        if img_height > base_height:
            img_width, img_height = (
                int(img_width * (base_height / img_height)),
                base_height,
            )
        img_left = int(base_left + 0.5 * (base_width - img_width))
        img_top = int(base_top + 0.5 * (base_height - img_height))
        slide.shapes.add_picture(
            image["imgfile"],
            left=img_left,
            top=img_top,
            width=img_width,
            height=img_height,
        )

    def put_question_on_slide(self, image, slide, q, allowbigimage=True):
        textbox = self.make_slide_layout(image, slide, allowbigimage=allowbigimage)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True
        self.set_question_number(slide, self.number)
        question = q["question"]
        handout = None
        if not image:
            handout, question = self._split_handout_from_text(question)
        question_text = self.pptx_process_text(question, image=image)
        question_size = self._get_font_size_for_text(
            "question", question_text, "question_size", 32
        )
        if handout:
            handout_cfg = self.c.get("handout", {})
            handout_text = self.pptx_process_text(
                handout, do_not_remove_accents=True
            )
            handout_p = self.init_paragraph(
                tf,
                size=self._get_handout_font_size_for_text(
                    handout_text, textbox, min_size=question_size
                ),
                line_spacing_key="handout",
            )
            self._set_paragraph_alignment(handout_p, handout_cfg.get("align"))
            self.pptx_format(
                handout_text,
                handout_p,
                tf,
                slide,
            )
            handout_p.space_after = PptxPt(self._get_handout_text_space_after())
            p = self._configure_paragraph(
                tf.add_paragraph(),
                size=question_size,
                line_spacing_key="question",
            )
        else:
            p = self.init_paragraph(
                tf,
                size=question_size,
                line_spacing_key="question",
            )
        self.pptx_format(question_text, p, tf, slide, blank_lines_between_items=True)
        self._custom_shrink_textbox(textbox)

    def recursive_join(self, s):
        if isinstance(s, str):
            return s
        if isinstance(s, list):
            return "\n".join(self.recursive_join(x) for x in s)

    def add_slide_with_handout(self, handout, number=None):
        slide = self.prs.slides.add_slide(self.QUESTION_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True
        if number is not None:
            self.set_question_number(slide, number)
        handout_cfg = self.c.get("handout", {})
        handout_text = self.pptx_process_text(handout, do_not_remove_accents=True)
        p = self.init_paragraph(
            tf,
            size=self._get_handout_font_size_for_text(handout_text, textbox),
            line_spacing_key="handout",
        )
        self._set_paragraph_alignment(p, handout_cfg.get("align"))
        self.pptx_format(handout_text, p, tf, slide)
        self._custom_shrink_textbox(textbox)

    def process_question_text(self, q):
        image = self._get_image_from_4s(q["question"])
        handout = self._get_handout_from_4s(q["question"])
        add_handout_on_separate_slide = self._add_handout_on_separate_slide()
        if image and add_handout_on_separate_slide:
            self.add_slide_with_image(image, number=self.number)
        elif handout and add_handout_on_separate_slide:
            self.add_slide_with_handout(handout, number=self.number)
        slide = self.prs.slides.add_slide(self.QUESTION_SLIDE)
        text_is_duplicated = bool(self.c.get("text_is_duplicated"))
        self.put_question_on_slide(
            image, slide, q, allowbigimage=not text_is_duplicated
        )
        if image and image["big"] and text_is_duplicated:
            self.add_slide_with_image(image, number=self.number)

    def _get_answer_grid_text(self, q, fields):
        result = []
        for field in fields:
            strip_brackets = field not in ("answer", "zachet")
            value = self.pptx_process_text(
                copy.deepcopy(q[field]), strip_brackets=strip_brackets
            )
            result.append(f"{self.get_label(q, field)}: {self._text_for_grid(value)}")
        return "\n".join(result)

    def add_answer_slide(self, q):
        slide = self.prs.slides.add_slide(self.ANSWER_SLIDE)
        if self.c.get("override_answer_caption"):
            self.set_question_number(slide, self.c["override_answer_caption"])
        else:
            self.set_question_number(slide, self.number)
        fields = ["answer"]
        if q.get("zachet") and self.c.get("add_zachet"):
            fields.append("zachet")
        if q.get("nezachet") and self.c.get("add_zachet"):
            fields.append("nezachet")
        if self.c["add_comment"] and "comment" in q:
            fields.append("comment")
        if self.c.get("add_source") and "source" in q:
            fields.append("source")
        if self.c.get("add_author") and "author" in q:
            fields.append("author")
        answer_size = self._get_font_size_for_text(
            "answer", self._get_answer_grid_text(q, fields), "answer_size", 32
        )
        textbox = None
        for field in fields:
            image = self._get_image_from_4s(q[field])
            if image:
                textbox = self.make_slide_layout(image, slide)
                break
        if not textbox:
            textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True

        p = self.init_paragraph(
            tf,
            size=answer_size,
            line_spacing_key="answer",
        )
        r = self.add_run(p, f"{self.get_label(q, 'answer')}: ")
        r.font.bold = True
        self.pptx_format(
            self.pptx_process_text(q["answer"], strip_brackets=False), p, tf, slide
        )
        if q.get("zachet") and self.c.get("add_zachet"):
            zachet_text = self.pptx_process_text(q["zachet"], strip_brackets=False)
            r = self.add_run(p, f"\n{self.get_label(q, 'zachet')}: ")
            r.font.bold = True
            self.pptx_format(zachet_text, p, tf, slide)
        if q.get("nezachet") and self.c.get("add_zachet"):
            nezachet_text = self.pptx_process_text(q["nezachet"])
            r = self.add_run(p, f"\n{self.get_label(q, 'nezachet')}: ")
            r.font.bold = True
            self.pptx_format(nezachet_text, p, tf, slide)
        if self.c["add_comment"] and "comment" in q:
            comment_text = self.pptx_process_text(q["comment"])
            r = self.add_run(p, f"\n{self.get_label(q, 'comment')}: ")
            r.font.bold = True
            self.pptx_format(comment_text, p, tf, slide)
        if self.c.get("add_source") and "source" in q:
            source_text = self.pptx_process_text(q["source"])
            r = self.add_run(p, f"\n{self.get_label(q, 'source')}: ")
            r.font.bold = True
            self.pptx_format(source_text, p, tf, slide)
        if self.c.get("add_author") and "author" in q:
            author_text = self.pptx_process_text(q["author"])
            r = self.add_run(p, f"\n{self.get_label(q, 'author')}: ")
            r.font.bold = True
            self.pptx_format(author_text, p, tf, slide)
        self._custom_shrink_textbox(textbox)

    def process_question(self, q):
        if "number" not in q:
            self.qcount += 1
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        self.number = str(self.qcount if "number" not in q else q["number"])

        if isinstance(q["question"], list):
            for i in range(len(q["question"][1])):
                qn = copy.deepcopy(q)
                qn["question"][1] = q["question"][1][: i + 1]
                self.process_question_text(qn)
        else:
            self.process_question_text(q)

        if self.c["add_plug"]:
            slide = self.prs.slides.add_slide(self.PLUG_SLIDE)
            self.set_question_number(slide, self.number)
        self.add_answer_slide(q)

    def init_paragraph(self, text_frame, size=None, line_spacing_key=None):
        p = text_frame.paragraphs[0]
        return self._configure_paragraph(
            p, size=size, line_spacing_key=line_spacing_key
        )

    def export(self, outfilename):
        self.outfilename = outfilename
        wd = os.getcwd()
        os.chdir(os.path.dirname(self.config_path))
        template = os.path.abspath(self.c["template_path"])
        os.chdir(wd)
        self.prs = Presentation(template)
        template_version = self.c.get("template_version", 1)
        layouts = self.prs.slide_layouts
        self.TITLE_SLIDE = layouts[self.c.get("title_slide_index", 0)]
        self.BLANK_SLIDE = layouts[self.c.get("blank_slide_index", 6)]
        if template_version >= 2:
            self.QUESTION_SLIDE = layouts[self.c.get("question_slide_index", 1)]
            self.ANSWER_SLIDE = layouts[self.c.get("answer_slide_index", 2)]
            self.PLUG_SLIDE = layouts[self.c.get("plug_slide_index", 3)]
        else:
            self.QUESTION_SLIDE = self.BLANK_SLIDE
            self.ANSWER_SLIDE = self.BLANK_SLIDE
            self.PLUG_SLIDE = self.BLANK_SLIDE
        self._prepare_service_slide_templates()
        self._add_service_slides("intro")
        self._processed_question_count = 0
        self._processed_tour_count = 0
        buffer = []
        for element in self.structure:
            if element[0] != "Question":
                buffer.append(element)
                continue
            if element[0] == "Question":
                if buffer:
                    if self._should_add_between_tours_slide(buffer):
                        self._add_service_slides("between_tours")
                    self.process_buffer(buffer)
                    buffer = []
                self.process_question(element[1])
                self._processed_question_count += 1
        self._add_service_slides("final")
        self._remove_service_slide_templates()
        self.prs.save(outfilename)
        if self.optimize_size:
            optimize_pptx_images(outfilename, quality=80)
        self.logger.info("Output: {}".format(outfilename))
