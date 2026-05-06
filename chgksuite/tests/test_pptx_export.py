from io import BytesIO
from pathlib import Path
import random
import urllib.parse
import zipfile

import toml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

from chgksuite.common import DefaultArgs
from chgksuite.composer.docx import _HYPERLINK_SAFE_CHARS
from chgksuite.composer.pptx import PptxExporter, optimize_pptx_images


ROOT = Path(__file__).resolve().parents[1]
RESOURCES = ROOT / "chgksuite" / "resources"


def _merge_config(base, updates):
    for key, value in updates.items():
        if value is None:
            base.pop(key, None)
            continue
        if isinstance(value, dict):
            base.setdefault(key, {})
            _merge_config(base[key], value)
        else:
            base[key] = value
    return base


def _config_path(tmp_path, updates=None):
    if not updates:
        return str(RESOURCES / "pptx_config.toml")
    config = toml.load(RESOURCES / "pptx_config.toml")
    config["template_path"] = str(RESOURCES / "template.pptx")
    _merge_config(config, updates)
    path = tmp_path / "pptx_config.toml"
    path.write_text(toml.dumps(config), encoding="utf8")
    return str(path)


def _pptx_args(tmp_path, font=None, config_updates=None, optimize_size=None):
    args = DefaultArgs(
        pptx_config=_config_path(tmp_path, config_updates),
        labels_file=str(RESOURCES / "labels_ru.toml"),
        regexes_file=str(RESOURCES / "regexes_ru.json"),
        language="ru",
        replace_no_break_spaces="on",
        replace_no_break_hyphens="on",
        font=font,
    )
    if optimize_size is not None:
        args.optimize_size = optimize_size
    return args


def _export_pptx(
    tmp_path,
    structure,
    font=None,
    config_updates=None,
    optimize_size=None,
):
    return Presentation(
        str(
            _export_pptx_path(
                tmp_path,
                structure,
                font=font,
                config_updates=config_updates,
                optimize_size=optimize_size,
            )
        )
    )


def _export_pptx_path(
    tmp_path,
    structure,
    font=None,
    config_updates=None,
    optimize_size=None,
):
    outfilename = tmp_path / "out.pptx"
    exporter = PptxExporter(
        structure,
        _pptx_args(
            tmp_path,
            font=font,
            config_updates=config_updates,
            optimize_size=optimize_size,
        ),
        {"tmp_dir": str(tmp_path), "targetdir": str(tmp_path)},
    )
    exporter.export(str(outfilename))
    return outfilename


def _write_noisy_png(path):
    rng = random.Random(0)
    image = Image.new("RGB", (180, 180))
    image.putdata(
        [
            (rng.randrange(256), rng.randrange(256), rng.randrange(256))
            for _ in range(180 * 180)
        ]
    )
    image.save(path)


def _soft_breaks_as_newlines(text):
    return text.replace("\v", "\n")


def _slide_text(slide):
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )


def _service_slide_template(tmp_path):
    prs = Presentation()

    def add_text_slide(text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(1), PptxInches(5), PptxInches(1)
        )
        textbox.text = text
        return slide

    add_text_slide("INTRO A")
    add_text_slide("INTRO B")
    add_text_slide("BREAK")
    final_slide = add_text_slide("FINAL")
    final_slide.shapes.add_picture(
        str(ROOT / "tests" / "test.jpg"),
        PptxInches(1),
        PptxInches(2),
        width=PptxInches(1),
    )
    add_text_slide("UNUSED PROTOTYPE")
    add_text_slide("TOUR ONE")
    add_text_slide("TOUR TWO")

    path = tmp_path / "service-template.pptx"
    prs.save(path)
    return path


def test_optimize_pptx_images_recompresses_png_as_jpeg(tmp_path):
    image_path = tmp_path / "source.png"
    _write_noisy_png(image_path)
    pptx_path = tmp_path / "image.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), PptxInches(1), PptxInches(1))
    prs.save(pptx_path)
    original_size = pptx_path.stat().st_size

    renamed_parts = optimize_pptx_images(pptx_path, quality=80)

    with zipfile.ZipFile(pptx_path) as pptx_file:
        names = pptx_file.namelist()
        content_types = pptx_file.read("[Content_Types].xml").decode("utf-8")
        rels = pptx_file.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
        image_names = [name for name in names if name.startswith("ppt/media/")]
        image_data = pptx_file.read(image_names[0])

    assert renamed_parts == {"ppt/media/image1.png": "ppt/media/image1.jpg"}
    assert image_names == ["ppt/media/image1.jpg"]
    assert image_data.startswith(b"\xff\xd8")
    assert pptx_path.stat().st_size < original_size
    assert 'Extension="jpg" ContentType="image/jpeg"' in content_types
    assert 'Target="../media/image1.jpg"' in rels


def test_optimize_pptx_images_preserves_transparent_png(tmp_path):
    image_path = tmp_path / "transparent.png"
    Image.new("RGBA", (120, 120), (255, 0, 0, 128)).save(
        image_path, format="PNG", compress_level=0
    )
    pptx_path = tmp_path / "transparent.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), PptxInches(1), PptxInches(1))
    prs.save(pptx_path)

    optimized_parts = optimize_pptx_images(pptx_path, quality=80)

    with zipfile.ZipFile(pptx_path) as pptx_file:
        names = pptx_file.namelist()
        rels = pptx_file.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
        image_names = [name for name in names if name.startswith("ppt/media/")]
        image_data = pptx_file.read(image_names[0])

    assert optimized_parts == {"ppt/media/image1.png": "ppt/media/image1.png"}
    assert image_names == ["ppt/media/image1.png"]
    assert not image_data.startswith(b"\xff\xd8")
    assert 'Target="../media/image1.png"' in rels
    with Image.open(BytesIO(image_data)) as image:
        assert image.convert("RGBA").getchannel("A").getextrema()[0] < 255
    Presentation(str(pptx_path))


def test_pptx_exporter_optimizes_size_by_default(tmp_path, monkeypatch):
    calls = []

    def fake_optimize(pptx_path, quality=80):
        calls.append((Path(pptx_path).name, quality))
        return {}

    monkeypatch.setattr("chgksuite.composer.pptx.optimize_pptx_images", fake_optimize)
    _export_pptx(
        tmp_path,
        [("Question", {"question": "Вопрос.", "answer": "Ответ."})],
    )

    assert calls == [("out.pptx", 80)]


def test_pptx_exporter_can_disable_size_optimization(tmp_path, monkeypatch):
    calls = []

    def fake_optimize(pptx_path, quality=80):
        calls.append((Path(pptx_path).name, quality))
        return {}

    monkeypatch.setattr("chgksuite.composer.pptx.optimize_pptx_images", fake_optimize)
    _export_pptx(
        tmp_path,
        [("Question", {"question": "Вопрос.", "answer": "Ответ."})],
        optimize_size="off",
    )

    assert calls == []


def test_pptx_hyperlinks_are_clickable_and_styled(tmp_path):
    url = "https://example.com/ик-с?q=тест"
    encoded_url = urllib.parse.quote(url, safe=_HYPERLINK_SAFE_CHARS)
    pptx_path = _export_pptx_path(
        tmp_path,
        [
            (
                "Question",
                {
                    "question": f"Ссылка: {url}.",
                    "answer": "Ответ.",
                },
            ),
        ],
    )
    prs = Presentation(str(pptx_path))

    link_runs = [
        run
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text == url
    ]

    assert len(link_runs) == 1
    assert link_runs[0].hyperlink.address == encoded_url
    assert link_runs[0].font.underline is True
    assert link_runs[0].font.color.rgb == RGBColor(0x05, 0x63, 0xC1)

    with zipfile.ZipFile(pptx_path) as pptx_file:
        slide_xml = "\n".join(
            pptx_file.read(name).decode("utf-8")
            for name in pptx_file.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        rels_xml = "\n".join(
            pptx_file.read(name).decode("utf-8")
            for name in pptx_file.namelist()
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels")
        )

    assert '<a:hlinkClick r:id="' in slide_xml
    assert 'u="sng"' in slide_xml
    assert 'val="0563C1"' in slide_xml
    assert 'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"' in rels_xml
    assert f'Target="{encoded_url}"' in rels_xml
    assert 'TargetMode="External"' in rels_xml


def test_pptx_can_disable_hyperlink_formatting(tmp_path):
    url = "https://example.com/ик-с?q=тест"
    encoded_url = urllib.parse.quote(url, safe=_HYPERLINK_SAFE_CHARS)
    pptx_path = _export_pptx_path(
        tmp_path,
        [
            (
                "Question",
                {
                    "question": f"Ссылка: {url}.",
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"format_links": False},
    )
    prs = Presentation(str(pptx_path))

    link_runs = [
        run
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text == url
    ]

    assert len(link_runs) == 1
    assert not link_runs[0].hyperlink.address
    assert link_runs[0].font.underline is not True

    with zipfile.ZipFile(pptx_path) as pptx_file:
        slide_xml = "\n".join(
            pptx_file.read(name).decode("utf-8")
            for name in pptx_file.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        rels_xml = "\n".join(
            pptx_file.read(name).decode("utf-8")
            for name in pptx_file.namelist()
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels")
        )

    assert '<a:hlinkClick r:id="' not in slide_xml
    assert encoded_url not in rels_xml


def test_pptx_block_metadata_hyperlinks_are_clickable_and_styled(tmp_path):
    url = "https://gotquestions.online"
    pptx_path = _export_pptx_path(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("section", "Тур 1"),
            ("meta", f"Срок незасветки пакета — до публикации на {url}."),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
    )
    prs = Presentation(str(pptx_path))

    link_runs = [
        run
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame")
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text == url
    ]

    assert len(link_runs) == 1
    assert link_runs[0].hyperlink.address == url
    assert link_runs[0].font.underline is True
    assert link_runs[0].font.color.rgb == RGBColor(0x05, 0x63, 0xC1)

    with zipfile.ZipFile(pptx_path) as pptx_file:
        rels_xml = "\n".join(
            pptx_file.read(name).decode("utf-8")
            for name in pptx_file.namelist()
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels")
        )

    assert (
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
        'relationships/hyperlink"'
    ) in rels_xml
    assert f'Target="{url}"' in rels_xml
    assert 'TargetMode="External"' in rels_xml


def test_pptx_export_preserves_zachet_brackets(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            (
                "Question",
                {
                    "question": "Вопрос [убрать].",
                    "answer": "Ответ [оставить].",
                    "zachet": "Зачет [оставить].",
                    "nezachet": "Незачет [убрать].",
                    "comment": "Комментарий [убрать].",
                },
            ),
        ],
    )

    text = "\n".join(_slide_text(slide) for slide in prs.slides)

    assert "Вопрос." in text
    assert "Ответ [оставить]." in text
    assert "Зачет [оставить]." in text
    assert "Незачет." in text
    assert "Комментарий." in text
    assert "[убрать]" not in text


def test_title_slide_uses_full_height_centered_textbox(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            (
                "heading",
                "Альфа. «Зеркало» первого игрового дня основной дисциплины СтудЧР-2026",
            ),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
    )

    title = prs.slides[0].shapes.title

    assert title.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert "<a:noAutofit" in title.element.txBody.xml
    assert "<a:normAutofit" not in title.element.txBody.xml
    assert title.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    assert round(title.left / 914400, 2) == 0.8
    assert round(title.top / 914400, 2) == 0.8
    assert round(title.height / 914400, 2) == 6.1
    assert round(title.width / 914400, 2) == 10.5
    assert [run.font.size.pt for run in title.text_frame.paragraphs[0].runs] == [60.0]
    assert not any(shape.name.startswith("Subtitle") for shape in prs.slides[0].shapes)


def test_default_pptx_template_uses_larger_tour_and_handout_fonts(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("section", "Тур 1"),
            (
                "Question",
                {
                    "question": "[Раздаточный материал: Текст раздатки]\nВопрос?",
                    "answer": "Ответ.",
                },
            ),
        ],
    )

    tour_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text.strip() == "Тур 1"
    )
    handout_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text.strip() == "Текст раздатки"
    )

    assert {
        run.font.size.pt
        for run in tour_shape.text_frame.paragraphs[0].runs
        if run.text.strip()
    } == {42.0}
    assert {
        run.font.size.pt
        for run in handout_shape.text_frame.paragraphs[0].runs
        if run.text.strip()
    } == {42.0}

    question_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and "Вопрос?" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]
    question_paragraph = question_shape.text_frame.paragraphs[1]

    assert handout_paragraph.text == "Текст раздатки"
    assert {
        run.font.size.pt for run in handout_paragraph.runs if run.text.strip()
    } == {42.0}
    assert question_paragraph.text == "Вопрос?"
    assert {
        run.font.size.pt for run in question_paragraph.runs if run.text.strip()
    } == {32.0}


def test_service_slides_are_inserted_from_template(tmp_path):
    template_path = _service_slide_template(tmp_path)

    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("section", "Тур 1"),
            ("Question", {"question": "Вопрос 1.", "answer": "Ответ 1."}),
            ("section", "Тур 2"),
            ("Question", {"question": "Вопрос 2.", "answer": "Ответ 2."}),
        ],
        config_updates={
            "template_path": str(template_path),
            "add_plug": False,
            "service_slides": {
                "intro": [0, 1],
                "between_tours": 2,
                "final": 3,
                "remove": [0, 1, 2, 3, 4, 5, 6],
                "skip_generated_title": True,
            },
        },
    )

    slide_texts = [_slide_text(slide) for slide in prs.slides]

    assert slide_texts[0] == "INTRO A"
    assert slide_texts[1] == "INTRO B"
    assert "Тестовый пакет" not in "\n".join(slide_texts)
    assert "UNUSED PROTOTYPE" not in "\n".join(slide_texts)
    assert slide_texts.count("BREAK") == 1
    assert slide_texts.index("BREAK") < next(
        index for index, text in enumerate(slide_texts) if "Тур 2" in text
    )
    assert "FINAL" in slide_texts[-1]
    assert any(shape.shape_type == 13 for shape in prs.slides[-1].shapes)


def test_numbered_tour_stubs_are_inserted_before_generated_tour_slides(tmp_path):
    template_path = _service_slide_template(tmp_path)

    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("editor", "Редактор: Тест."),
            ("section", "Тур 1"),
            ("Question", {"question": "Вопрос 1.", "answer": "Ответ 1."}),
            ("section", "Тур 2"),
            ("Question", {"question": "Вопрос 2.", "answer": "Ответ 2."}),
        ],
        config_updates={
            "template_path": str(template_path),
            "add_plug": False,
            "service_slides": {
                "numbered_tours_stubs": [5, 6],
                "remove": [0, 1, 2, 3, 4, 5, 6],
                "skip_generated_title": True,
            },
        },
    )

    slide_texts = [_slide_text(slide) for slide in prs.slides]

    assert slide_texts.index("Редактор: Тест.") < slide_texts.index("TOUR ONE")
    assert slide_texts.index("TOUR ONE") < next(
        index for index, text in enumerate(slide_texts) if "Тур 1" in text
    )
    assert slide_texts.index("TOUR TWO") < next(
        index for index, text in enumerate(slide_texts) if "Тур 2" in text
    )


def test_pptx_textboxes_shrink_text_and_stamp_run_sizes(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("section", "Тур 2"),
            (
                "editor",
                "Редакторы: Александр Сновский, Артём Сапожников, "
                "Александр Зайцев, Ольга Шиншинова, Мария Аристова.",
            ),
            (
                "meta",
                "Редакторы благодарят за тестирование вопросов и очень ценные "
                "замечания: Андрея Багдуева, Дмитрия Батова, Марию Богуш, "
                "Виталия Буковского, Михаила Гриценко, Дарью Жукову, "
                "Фёдора Журавлёва, Григория Зырянова, Алексея Ковбу, "
                "Евгению Колпащикову, Юстину Кустовскую, Маргариту Лузину, "
                "Михаила Малкина, Веру Монину, Ерлана Мухамеджанова, "
                "Илью Орлова, Никиту Пеговса, Кирилла Платонова, "
                "Юрия Разумова, Дмитрия Селянина, Дмитрия Слободянюка, "
                "Наиля Фарукшина, Ксению и Эдуарда Шагалов и Максима Шиловского.",
            ),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
    )

    textboxes = [
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Редакторы:" in shape.text
    ]
    assert len(textboxes) == 1
    textbox = textboxes[0]

    assert textbox.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert "<a:noAutofit" in textbox.element.txBody.xml
    assert "<a:normAutofit" not in textbox.element.txBody.xml
    assert "<a:lnSpc>" not in textbox.element.txBody.xml
    assert "<a:spAutoFit" not in textbox.element.txBody.xml

    runs = [
        run
        for paragraph in textbox.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    ]
    run_sizes = [run.font.size.pt for run in runs]
    assert run_sizes
    assert {run.font.size.pt for run in runs if "Тур 2" in run.text} == {34.0}
    assert {run.font.size.pt for run in runs if "Тур 2" not in run.text} == {24.0}
    assert {paragraph.line_spacing for paragraph in textbox.text_frame.paragraphs} == {
        None
    }


def test_fixed_line_spacing_keeps_config_spacing_with_shrink_fit(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
        config_updates={
            "font": {
                "question_size": 24,
                "answer_size": 20,
                "number_size": 26,
                "fixed_line_spacing_question": 24,
                "fixed_line_spacing_answer": 20,
            },
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Вопрос." in shape.text
    )
    assert question_shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert "<a:noAutofit" in question_shape.element.txBody.xml
    assert "<a:normAutofit" not in question_shape.element.txBody.xml
    assert '<a:spcPts val="2400"/>' in question_shape.element.txBody.xml
    assert {
        paragraph.line_spacing.pt for paragraph in question_shape.text_frame.paragraphs
    } == {24.0}

    number_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and shape.text == "1"
    )
    assert number_shape.text_frame.paragraphs[0].line_spacing is None

    answer_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and "Ответ." in shape.text
    )
    assert '<a:spcPts val="2000"/>' in answer_shape.element.txBody.xml
    assert {
        paragraph.line_spacing.pt for paragraph in answer_shape.text_frame.paragraphs
    } == {20.0}


def test_line_spacing_multiplier_sets_percent_spacing(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
        config_updates={
            "font": {
                "question_size": 24,
                "answer_size": 20,
                "line_spacing_multiplier": 1.2,
            },
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Вопрос." in shape.text
    )
    assert "<a:normAutofit" not in question_shape.element.txBody.xml
    assert '<a:spcPct val="120000"/>' in question_shape.element.txBody.xml
    assert {
        paragraph.line_spacing for paragraph in question_shape.text_frame.paragraphs
    } == {1.2}

    answer_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and "Ответ." in shape.text
    )
    assert '<a:spcPct val="120000"/>' in answer_shape.element.txBody.xml
    assert {
        paragraph.line_spacing for paragraph in answer_shape.text_frame.paragraphs
    } == {1.2}


def test_disable_shrink_fit_uses_no_autofit_and_config_font_sizes(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
        config_updates={
            "disable_shrink_fit": True,
            "font": {"question_size": 24, "answer_size": 20, "number_size": 26},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Вопрос." in shape.text
    )
    assert question_shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert "<a:noAutofit" in question_shape.element.txBody.xml
    assert "<a:normAutofit" not in question_shape.element.txBody.xml
    assert {
        run.font.size.pt
        for paragraph in question_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {24.0}

    number_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and shape.text == "1"
    )
    assert number_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 26.0

    answer_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and "Ответ." in shape.text
    )
    assert {
        run.font.size.pt
        for paragraph in answer_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {20.0}


def test_text_size_grid_sets_question_and_answer_sizes(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": "Очень длинный текст вопроса для проверки сетки.",
                    "answer": "Очень длинный текст ответа для проверки сетки.",
                    "comment": "Комментарий тоже участвует в размере ответа.",
                },
            ),
        ],
        config_updates={
            "disable_shrink_fit": True,
            "text_size_grid": {
                "question_elements": [
                    {"length": 10, "size": 24},
                    {"length": 1000, "size": 18},
                ],
                "answer_elements": [
                    {"length": 10, "size": 20},
                    {"length": 1000, "size": 16},
                ],
                "smallest": 14,
            },
            "font": {"question_size": 24, "answer_size": 20},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Очень длинный текст вопроса" in shape.text
    )
    assert question_shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert {
        run.font.size.pt
        for paragraph in question_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {18.0}

    answer_shape = next(
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and "Очень длинный текст ответа" in shape.text
    )
    assert {
        run.font.size.pt
        for paragraph in answer_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {16.0}


def test_number_textbox_font_size_overrides_font_number_size(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
        config_updates={
            "font": {"number_size": 26},
            "number_textbox": {"font_size": 28},
        },
    )

    number_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and shape.text == "1"
    )

    assert number_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 28.0


def test_pptx_font_override_replaces_config_font(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
        font="Times New Roman",
    )

    title = prs.slides[0].shapes.title
    assert title.text_frame.paragraphs[0].runs[0].font.name == "Times New Roman"

    question_runs = [
        run
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Вопрос." in shape.text
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    ]
    assert question_runs
    assert {run.font.name for run in question_runs} == {"Times New Roman"}


def test_caps_question_number_formats_zero_with_label(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {"number": 0, "question": "Разминка.", "answer": "Ответ."},
            ),
        ],
        config_updates={"question_number_format": "caps"},
    )

    number_texts = [
        shape.text
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame")
    ]

    assert "ВОПРОС 0" in number_texts


def test_douplet_list_gets_extra_break_before_numbering(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": [
                        "Дуплет. Два вопроса по 30 секунд каждый.",
                        [
                            "Первый подвопрос.",
                            "Второй подвопрос.",
                        ],
                    ],
                    "answer": ["Первый.", "Второй."],
                },
            ),
        ],
    )

    question_text = next(
        shape.text
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Дуплет." in shape.text
    )

    assert (
        "каждый.\n\n1. Первый подвопрос."
        in _soft_breaks_as_newlines(question_text)
    )


def test_douplet_list_numbering_style_is_configurable(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": [
                        "Блиц.",
                        [
                            "Первый подвопрос.",
                            "Второй подвопрос.",
                        ],
                    ],
                    "answer": ["Первый.", "Второй."],
                },
            ),
        ],
        config_updates={"list": {"numbering_style": "a)"}},
    )

    question_text = next(
        shape.text
        for shape in prs.slides[2].shapes
        if hasattr(shape, "text_frame") and "Блиц." in shape.text
    )

    assert (
        "Блиц.\n\na) Первый подвопрос.\n\nb) Второй подвопрос."
        in _soft_breaks_as_newlines(question_text)
    )


def test_source_lists_stay_compact_on_answer_slide(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": "Вопрос.",
                    "answer": "Ответ.",
                    "source": [
                        "Первый источник.",
                        "Второй источник.",
                    ],
                },
            ),
        ],
        config_updates={"add_source": True, "list": {"numbering_style": "1)"}},
    )

    answer_text = next(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and "Источники:" in shape.text
    )
    answer_text = _soft_breaks_as_newlines(answer_text)

    assert "\n1) Первый источник.\n2) Второй источник." in answer_text
    assert "\n\n2) Второй источник." not in answer_text


def test_handout_slide_uses_handout_config_and_hides_label_by_default(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": "[Раздаточный материал: Текст раздатки]\nВопрос?",
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"handout": {"font_size": 24, "align": "center"}},
    )

    handout_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Текст раздатки" in shape.text
    )
    paragraph = handout_shape.text_frame.paragraphs[0]

    assert handout_shape.text == "Текст раздатки"
    assert paragraph.alignment == PP_ALIGN.CENTER
    assert {run.font.size.pt for run in paragraph.runs if run.text.strip()} == {24.0}


def test_handout_label_can_be_included_from_config(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": "[Раздаточный материал: Текст раздатки]\nВопрос?",
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"handout": {"include_label": True}},
    )

    handout_text = next(
        shape.text
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Текст раздатки" in shape.text
    )

    assert handout_text == "[Раздаточный материал: Текст раздатки]"


def test_inline_handout_uses_handout_config_when_not_separate_slide(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        "[Раздаточный материал:\n"
                        "Klemperer\n"
                        "]\n"
                        "Чтобы спастись от гестапо, супруги Клемперер решили "
                        "подделать свои документы."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_handout_on_separate_slide": False,
            "force_text_size_question": 24,
            "font": {"default_size": None, "question_size": None},
            "handout": {"include_label": False, "font_size": 32, "align": "center"},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Klemperer" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]
    question_paragraph = question_shape.text_frame.paragraphs[1]

    assert handout_paragraph.text == "Klemperer"
    assert handout_paragraph.alignment == PP_ALIGN.CENTER
    assert {
        run.font.size.pt for run in handout_paragraph.runs if run.text.strip()
    } == {32.0}
    assert handout_paragraph.space_after.pt == 72.0
    assert question_paragraph.text.startswith("Чтобы спастись")
    assert {
        run.font.size.pt for run in question_paragraph.runs if run.text.strip()
    } == {24.0}


def test_inline_handout_line_uses_larger_default_handout_font(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        "Раздаточный материал: Klemperer\n"
                        "Чтобы спастись от гестапо, супруги Клемперер решили "
                        "подделать свои документы."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_handout_on_separate_slide": False,
            "force_text_size_question": 24,
            "font": {"default_size": None, "question_size": None},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Klemperer" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]
    question_paragraph = question_shape.text_frame.paragraphs[1]

    assert handout_paragraph.text == "Klemperer"
    assert {
        run.font.size.pt for run in handout_paragraph.runs if run.text.strip()
    } == {42.0}
    assert question_paragraph.text.startswith("Чтобы спастись")
    assert {
        run.font.size.pt for run in question_paragraph.runs if run.text.strip()
    } == {24.0}


def test_inline_handout_font_size_shrinks_to_preserve_source_lines(
    tmp_path, monkeypatch
):
    monkeypatch.setattr(
        PptxExporter, "_get_measurement_font_path", lambda self: "fake-font"
    )
    monkeypatch.setattr(
        PptxExporter,
        "_measure_text_width_px",
        lambda self, font_path, text, size: len(text) * size,
    )
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        "[Раздаточный материал: abcdefghijklmnopqrstuvwxy]\n"
                        "Здесь начинается текст вопроса"
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "force_text_size_question": 32,
            "font": {"default_size": None, "question_size": None},
            "handout": {"include_label": False, "font_size": 42},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[2].shapes
        if hasattr(shape, "text_frame") and "abcdefghijklmnopqrstuvwxy" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]
    question_paragraph = question_shape.text_frame.paragraphs[1]

    assert handout_paragraph.text == "abcdefghijklmnopqrstuvwxy"
    assert {
        run.font.size.pt for run in handout_paragraph.runs if run.text.strip()
    } == {39.0}
    assert {
        run.font.size.pt for run in question_paragraph.runs if run.text.strip()
    } == {32.0}


def test_inline_handout_text_spacing_can_be_configured(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        "[Раздаточный материал: Klemperer]\n"
                        "Чтобы спастись от гестапо, супруги Клемперер решили "
                        "подделать свои документы."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_handout_on_separate_slide": False,
            "handout": {
                "include_label": False,
                "font_size": 32,
                "space_after": 42,
                "text_space_after": 7,
            },
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Klemperer" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]

    assert handout_paragraph.space_after.pt == 7.0


def test_multiline_inline_handout_uses_soft_breaks_without_extra_paragraphs(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        "[Раздаточный материал:\n"
                        "привет\n"
                        "сегодня\n"
                        "такая\n"
                        "раздатка\n"
                        "]\n"
                        "Здесь начинается текст вопроса"
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_handout_on_separate_slide": False,
            "force_text_size_question": 24,
            "font": {"default_size": None, "question_size": None},
            "handout": {"include_label": False, "font_size": 32, "align": "center"},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "привет" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]

    assert len(question_shape.text_frame.paragraphs) == 2
    assert _soft_breaks_as_newlines(handout_paragraph.text) == (
        "привет\nсегодня\nтакая\nраздатка"
    )
    assert handout_paragraph._p.xml.count("<a:br") == 3
    assert "<a:t>привет\nсегодня" not in handout_paragraph._p.xml


def test_handout_image_scale_and_spacing_apply_to_question_slide(tmp_path):
    image_path = ROOT / "tests" / "test.jpg"
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        f"[Раздаточный материал: (img w=3in {image_path})]\n"
                        "Перед вами картинка. Назовите ее."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"handout": {"image_scale": 1.5, "space_after": 24}},
    )

    question_slide = next(
        slide
        for slide in prs.slides
        if any(
            hasattr(shape, "text_frame")
            and "Перед вами картинка" in shape.text.replace("\xa0", " ")
            for shape in slide.shapes
        )
    )
    picture = next(shape for shape in question_slide.shapes if shape.shape_type == 13)
    question_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame")
        and "Перед вами картинка" in shape.text.replace("\xa0", " ")
    )

    assert picture.width == PptxInches(4.5)
    assert question_shape.top == picture.top + picture.height + PptxPt(24)


def test_overlay_image_and_text_keeps_textbox_on_base_rect(tmp_path):
    image_path = ROOT / "tests" / "test.jpg"
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        f"(img w=3in {image_path})\n"
                        "Перед вами картинка. Назовите ее."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"overlay_image_and_text": True},
    )

    question_slide = next(
        slide
        for slide in prs.slides
        if any(
            hasattr(shape, "text_frame")
            and "Перед вами картинка" in shape.text.replace("\xa0", " ")
            for shape in slide.shapes
        )
    )
    picture = next(shape for shape in question_slide.shapes if shape.shape_type == 13)
    question_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame")
        and "Перед вами картинка" in shape.text.replace("\xa0", " ")
    )

    assert picture.left == PptxInches(0.8)
    assert picture.top == PptxInches(0.8)
    assert question_shape.left == PptxInches(0.8)
    assert question_shape.top == PptxInches(0.8)
    assert question_shape.left < picture.left + picture.width
    assert question_shape.top < picture.top + picture.height


def test_legacy_pptx_config_sizes_and_disabled_autolayout_do_not_overlap(tmp_path):
    image_path = ROOT / "tests" / "test.jpg"
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        f"[Раздаточный материал: (img {image_path})]\n"
                        "Перед вами картинка. Назовите ее."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_plug": False,
            "add_handout_on_separate_slide": False,
            "disable_autolayout": True,
            "force_text_size_question": 24,
            "force_text_size_answer": 20,
            "text_size_grid": {"default": 24},
            "number_textbox": {"font_size": 28},
            "font": {"default_size": None, "question_size": None, "answer_size": None},
        },
    )

    question_slide = prs.slides[1]
    picture = next(shape for shape in question_slide.shapes if shape.shape_type == 13)
    question_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame")
        and "Перед вами картинка" in shape.text.replace("\xa0", " ")
    )
    number_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame") and shape.text == "1"
    )

    assert question_shape.top >= picture.top + picture.height
    assert {
        run.font.size.pt
        for paragraph in question_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {24.0}
    assert number_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 28.0

    answer_shape = next(
        shape
        for shape in prs.slides[2].shapes
        if hasattr(shape, "text_frame") and "Ответ." in shape.text
    )
    assert {
        run.font.size.pt
        for paragraph in answer_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {20.0}
    assert 'lnSpcReduction="0"' not in answer_shape.element.txBody.xml
